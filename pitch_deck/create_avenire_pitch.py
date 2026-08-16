from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pitch_deck" / "avenire_pitch_deck.pptx"
LOGO = ROOT / "apps/web/public/branding/avenire-logo-full.png"
OG_UI = ROOT / "apps/web/public/og/avenire.png"
PRICING = Path("/home/apollo/.t3/userdata/attachments/6269efa1-cd2e-49b8-a863-0dbdb0483799-19278aa2-8aa4-4274-a790-d9bf9881845f.png")


# Avenire dark-mode tokens from packages/ui/src/styles.css.
BG = "141414"
SURFACE = "181818"
SURFACE_2 = "1D1D1D"
SURFACE_3 = "222222"
WHITE = "F4F4F4"
TEXT = "E4E4E4"
MUTED = "9C9C9C"
FAINT = "6E6E6E"
BLUE = "ABC4FF"
TEAL = "78C4D5"
ORANGE = "F0A873"
PINK = "ED84BD"
PURPLE = "A99AD1"
YELLOW = "FFDC49"
GREEN = "4DB39F"
RED = "FC6B83"
LINE = "343434"

FONT = "Inter"
MONO = "Cascadia Mono"

W = 13.333
H = 7.5
M = 0.65


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_fill(shape, color: str, transparency: int = 0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def no_fill(shape):
    shape.fill.background()


def set_line(shape, color: str = LINE, width: float = 1.0, transparency: int = 0):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, fill=SURFACE, line=None, radius=False, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    set_fill(shape, fill, transparency)
    if line:
        set_line(shape, line)
    else:
        no_line(shape)
    return shape


def circle(slide, x, y, d, fill=BLUE, line=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    set_fill(shape, fill, transparency)
    if line:
        set_line(shape, line)
    else:
        no_line(shape)
    return shape


def line(slide, x1, y1, x2, y2, color=LINE, width=1.4, dash=None):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    set_line(shape, color, width)
    if dash:
        shape.line.dash_style = dash
    return shape


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    size=18,
    color=TEXT,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    italic=False,
    margin=0,
    tracking=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    if tracking is not None:
        run.font.kerning = Pt(tracking)
    return box


def rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for value, color, bold in runs:
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
    return box


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def brand(slide, index=None, dark=True):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(M), Inches(0.28), height=Inches(0.23))
    text(slide, "Avenire", M + 0.31, 0.26, 0.9, 0.26, size=11, color=TEXT, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if index is not None:
        text(slide, f"{index:02d}", W - M - 0.35, 0.29, 0.35, 0.2, size=9, color=FAINT, font=MONO, align=PP_ALIGN.RIGHT)


def base_slide(prs, index=None, light=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = "F7F7F5" if light else BG
    rect(slide, 0, 0, W, H, fill=bg)
    if not light:
        brand(slide, index)
    return slide


def title_block(slide, kicker, title_value, sub=None, y=0.82, light=False):
    fg = "1A1A1A" if light else TEXT
    muted = "6C6C6C" if light else MUTED
    if kicker:
        text(slide, kicker.upper(), M, y, 5.8, 0.22, size=10, color=BLUE if not light else "5474B8", bold=True, font=MONO)
    text(slide, title_value, M, y + 0.28, 11.8, 0.7, size=30, color=fg, bold=True)
    if sub:
        text(slide, sub, M, y + 1.03, 10.6, 0.42, size=14, color=muted)


def pill(slide, label, x, y, w, color=BLUE, fill=SURFACE_3, size=10):
    # Plain section marker: avoid turning every small label into a badge.
    line(slide, x, y + 0.28, x + w, y + 0.28, color, 1.0)
    text(slide, label.upper(), x, y + 0.01, w, 0.18, size=size, color=color, bold=True, font=MONO, align=PP_ALIGN.CENTER)


def node(slide, label, x, y, w, h=0.55, fill=SURFACE_2, accent=BLUE, label_size=13, sub=None):
    rect(slide, x, y, w, h, fill=fill, line=accent, radius=False)
    text(slide, label, x + 0.12, y + 0.08, w - 0.24, 0.24, size=label_size, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        text(slide, sub, x + 0.12, y + 0.31, w - 0.24, 0.16, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)


def arrow_down(slide, x, y, h=0.35, color=BLUE):
    line(slide, x, y, x, y + h - 0.08, color, 1.8)
    tri = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x - 0.07), Inches(y + h - 0.17), Inches(0.14), Inches(0.17))
    set_fill(tri, color)
    no_line(tri)


def chevron(slide, x, y, w=0.26, h=0.36, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, color)
    no_line(shape)
    return shape


def mini_icon(slide, kind, x, y, color=BLUE):
    # Small editable glyphs built from PowerPoint shapes.
    if kind == "pdf":
        rect(slide, x, y, 0.28, 0.34, fill=SURFACE_3, line=color, radius=True)
        line(slide, x + 0.06, y + 0.11, x + 0.22, y + 0.11, color, 1)
        line(slide, x + 0.06, y + 0.18, x + 0.19, y + 0.18, color, 1)
    elif kind == "video":
        rect(slide, x, y, 0.28, 0.34, fill=SURFACE_3, line=color, radius=True)
        tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(x + 0.09), Inches(y + 0.09), Inches(0.12), Inches(0.15))
        set_fill(tri, color)
        no_line(tri)
    elif kind == "chat":
        rect(slide, x, y, 0.32, 0.25, fill=SURFACE_3, line=color, radius=True)
        tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(x + 0.04), Inches(y + 0.19), Inches(0.09), Inches(0.08))
        set_fill(tri, color)
        no_line(tri)
    else:
        circle(slide, x, y, 0.28, fill=SURFACE_3, line=color)


def add_footer(slide, label="Avenire / learning that remembers"):
    line(slide, M, H - 0.43, W - M, H - 0.43, LINE, 0.6)
    text(slide, label, M, H - 0.31, 3.8, 0.16, size=8.5, color=FAINT, font=MONO)


def add_bar(slide, x, y, max_w, value, color, h=0.13, radius=True):
    rect(slide, x, y, max_w, h, fill=SURFACE_3, radius=radius)
    rect(slide, x, y, max_w * value, h, fill=color, radius=radius)


def add_table(slide, rows, x, y, w, h, col_widths=None, header=True, font_size=11):
    tbl_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = tbl_shape.table
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
    row_h = h / len(rows)
    for r, row in enumerate(rows):
        table.rows[r].height = Inches(row_h)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(SURFACE_3 if r == 0 and header else (SURFACE_2 if r % 2 else SURFACE))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size if r else font_size - 0.5)
                    run.font.bold = r == 0 and header
                    run.font.color.rgb = rgb(BLUE if r == 0 and header else TEXT)
            # Thin, quiet borders.
            tc_pr = cell._tc.get_or_add_tcPr()
            for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                ln = tc_pr.find(edge, tc_pr.nsmap)
                if ln is None:
                    from pptx.oxml.xmlchemy import OxmlElement
                    ln = OxmlElement(edge)
                    tc_pr.append(ln)
                ln.set("w", "6000")
                solid = ln.find("a:solidFill", ln.nsmap)
                if solid is None:
                    from pptx.oxml.xmlchemy import OxmlElement
                    solid = OxmlElement("a:solidFill")
                    ln.append(solid)
                srgb = solid.find("a:srgbClr", solid.nsmap)
                if srgb is None:
                    from pptx.oxml.xmlchemy import OxmlElement
                    srgb = OxmlElement("a:srgbClr")
                    solid.append(srgb)
                srgb.set("val", LINE)
    return tbl_shape


def cover(prs):
    slide = base_slide(prs, None)
    # Very subtle knowledge convergence: editable lines and nodes, low contrast.
    pts = [(7.6, 1.1), (9.5, 2.1), (11.2, 1.35), (10.2, 3.0), (8.4, 3.8), (6.8, 3.05), (9.0, 4.45)]
    for i in range(len(pts) - 1):
        line(slide, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], PURPLE, 1.0)
    line(slide, 7.6, 1.1, 10.2, 3.0, TEAL, 1.0)
    line(slide, 9.5, 2.1, 8.4, 3.8, ORANGE, 1.0)
    for i, (x, y) in enumerate(pts):
        circle(slide, x - 0.04, y - 0.04, 0.08, fill=BLUE if i in (1, 3, 6) else SURFACE_3, line=PURPLE)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(M), Inches(1.2), height=Inches(0.58))
    text(slide, "Avenire", M + 0.78, 1.2, 2.2, 0.62, size=31, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    text(slide, "Learning that remembers.", M, 2.25, 8.8, 0.66, size=32, color=TEXT, bold=True)
    text(slide, "AI-native learning for students.", M, 3.05, 6.4, 0.4, size=16, color=MUTED)
    text(slide, "Abhiram — Founder | BITS Pilani, Hyderabad", M, 6.72, 6.5, 0.2, size=10, color=FAINT, font=MONO)
    add_notes(slide, "Hi, I'm Abhiram, and I'm building Avenire.\n\nAvenire is an AI-native learning platform built around one simple idea: learning should remember.\n\nToday, students have access to more information and AI tools than ever before, but the actual learning experience is still fragmented.")


def problem(prs):
    slide = base_slide(prs, 2)
    title_block(slide, "The problem", "Learning is fragmented.", "Context is scattered across the tools students already use.")
    cx, cy = 6.55, 3.35
    # Native network diagram.
    for x, y, label, accent, kind in [
        (2.15, 2.3, "PDF", ORANGE, "pdf"),
        (6.0, 1.4, "YouTube", PINK, "video"),
        (10.2, 2.3, "ChatGPT", TEAL, "chat"),
        (2.55, 4.55, "Notes", PURPLE, "notes"),
        (9.85, 4.55, "Flashcards", YELLOW, "cards"),
    ]:
        line(slide, cx, cy, x + 0.6, y + 0.25, accent, 1.4)
        node(slide, label, x, y, 1.6, 0.55, fill=SURFACE_2, accent=accent, label_size=12)
        mini_icon(slide, kind, x + 0.12, y + 0.12, accent)
    circle(slide, cx - 0.45, cy - 0.45, 0.9, fill=BLUE, line=BLUE)
    text(slide, "student", cx - 0.43, cy - 0.02, 0.86, 0.2, size=10, color=BG, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "Every tool knows a piece.", M, 6.05, 5.2, 0.34, size=19, color=TEXT, bold=True)
    text(slide, "None knows the learner.", 7.0, 6.05, 5.5, 0.34, size=19, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide)
    add_notes(slide, "A student might learn a topic from a lecture, look something up in a PDF, ask ChatGPT a question, make notes somewhere else, and later use a completely different app for flashcards.\n\nEach of these tools solves one part of learning.\n\nBut none of them really knows what the student has learned, what they're struggling with, or what they should do next.\n\nSo every session starts with fragmented context.")


def insight(prs):
    slide = base_slide(prs, 3)
    title_block(slide, "The insight", "AI shouldn't just answer questions.", "It should build understanding.")
    labels = [("ASK", BLUE), ("UNDERSTAND", TEAL), ("PRACTICE", ORANGE), ("REMEMBER", PURPLE)]
    x = 4.55
    for i, (label, accent) in enumerate(labels):
        y = 2.1 + i * 0.82
        line(slide, x, y + 0.5, x + 4.25, y + 0.5, accent, 1.4)
        text(slide, label, x, y + 0.13, 4.25, 0.2, size=14, color=accent, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            arrow_down(slide, x + 2.12, y + 0.52, 0.27, accent)
    text(slide, "Every interaction should make the next one better.", M, 6.15, 11.9, 0.35, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "This led us to a different way of thinking about AI in education.\n\nWe don't think the goal is to build another chatbot that gives students answers.\n\nThe goal is to build a system where every interaction contributes to an evolving understanding of the student.\n\nIf I ask a question today, that shouldn't disappear when the conversation ends. It should become part of what Avenire knows about my learning.")


def product(prs):
    slide = base_slide(prs, 4)
    title_block(slide, "The product", "Meet Avenire.", "Conversational Learning: ask, learn, create, remember.")
    # Actual Avenire UI asset inside an editable frame.
    rect(slide, 0.75, 2.0, 8.1, 4.62, fill="0B0B0B", line=LINE, radius=False)
    if OG_UI.exists():
        slide.shapes.add_picture(str(OG_UI), Inches(0.88), Inches(2.12), width=Inches(7.84), height=Inches(4.12))
    line(slide, 9.1, 2.18, 9.1, 6.42, LINE, 1.0)
    pill(slide, "A connected workflow", 9.67, 2.3, 2.55, color=TEAL, fill=SURFACE_3, size=8)
    steps = [("Question", BLUE), ("Explanation", TEAL), ("Interactive concept", ORANGE), ("Notes", PURPLE), ("Practice", PINK)]
    for i, (label, accent) in enumerate(steps):
        y = 2.95 + i * 0.56
        circle(slide, 9.72, y + 0.04, 0.18, fill=accent)
        text(slide, label, 10.08, y, 2.1, 0.25, size=12, color=TEXT, bold=i == 0)
        if i < len(steps) - 1:
            arrow_down(slide, 9.81, y + 0.23, 0.32, LINE)
    text(slide, "The product is not a collection of features.\nIt is one learning workflow.", 9.67, 5.87, 2.55, 0.5, size=11, color=MUTED)
    add_footer(slide)
    add_notes(slide, "This is Avenire.\n\nWe call the experience Conversational Learning.\n\nA student can ask a question, explore an explanation, interact with a visualization, create notes, generate practice material and eventually review that concept again.\n\nThese aren't separate workflows. They're connected parts of the same learning process.")


def loop(prs):
    slide = base_slide(prs, 5)
    title_block(slide, "The difference", "The conversation doesn't end when the chat ends.", "A generic assistant returns an answer. Avenire compounds learning.")
    labels = [("Conversation", BLUE), ("Knowledge", TEAL), ("Notes", ORANGE), ("Practice", PINK), ("Memory", PURPLE)]
    xs = [1.15, 3.5, 5.85, 8.2, 10.55]
    y = 3.0
    for i, ((label, accent), x) in enumerate(zip(labels, xs)):
        circle(slide, x, y, 0.72, fill=accent)
        text(slide, str(i + 1), x, y + 0.18, 0.72, 0.22, size=12, color=BG, bold=True, align=PP_ALIGN.CENTER, font=MONO)
        text(slide, label, x - 0.35, 4.0, 1.42, 0.25, size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            line(slide, x + 0.8, y + 0.36, xs[i + 1] - 0.1, y + 0.36, accent, 1.6)
            chevron(slide, xs[i + 1] - 0.22, y + 0.21, 0.22, 0.3, accent)
    # Return loop using editable lines/arrowheads.
    line(slide, 11.0, 3.0, 11.0, 2.3, PURPLE, 1.4)
    line(slide, 11.0, 2.3, 1.5, 2.3, PURPLE, 1.4)
    line(slide, 1.5, 2.3, 1.5, 2.95, PURPLE, 1.4)
    chevron(slide, 1.36, 2.84, 0.28, 0.3, PURPLE)
    text(slide, "One continuous learning loop.", M, 5.55, 11.9, 0.45, size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "This is the fundamental difference between Avenire and a generic AI assistant.\n\nWith a normal chatbot, I ask a question, get an answer, and move on.\n\nWith Avenire, that interaction can become part of my learning system.\n\nA conversation can become knowledge. Knowledge can become practice. Practice can become memory.\n\nAnd that memory improves future interactions.")


def technology(prs):
    slide = base_slide(prs, 6)
    title_block(slide, "The technology", "A learning intelligence layer underneath.", "Multimodal material in. Contextual, evidence-aware learning out.")
    # Open material column; the pipeline itself provides the structure.
    line(slide, 0.8, 2.78, 3.55, 2.78, LINE, 1.0)
    text(slide, "YOUR MATERIAL", 1.1, 2.45, 2.4, 0.25, size=11, color=BLUE, bold=True, font=MONO)
    materials = [("PDFs", ORANGE), ("Videos", PINK), ("Images", TEAL), ("Web", PURPLE), ("Notes", YELLOW), ("Chats", BLUE)]
    for i, (label, accent) in enumerate(materials):
        col = i % 2
        row = i // 2
        x = 1.08 + col * 1.37
        y = 3.1 + row * 0.7
        circle(slide, x, y + 0.02, 0.16, fill=accent)
        text(slide, label, x + 0.24, y, 0.95, 0.22, size=11, color=TEXT)
    text(slide, "Many formats.\nOne learning context.", 1.08, 5.0, 2.2, 0.45, size=13, color=MUTED)
    # Middle vertical stack.
    steps = [("UNDERSTAND", BLUE), ("CONNECT", TEAL), ("RETRIEVE", ORANGE), ("REASON", PINK)]
    x = 5.0
    for i, (label, accent) in enumerate(steps):
        y = 2.18 + i * 0.72
        node(slide, label, x, y, 2.35, 0.48, fill=SURFACE_2, accent=accent, label_size=11)
        if i == 0:
            line(slide, 3.9, 3.85, x, y + 0.24, BLUE, 1.6)
        if i < len(steps) - 1:
            arrow_down(slide, x + 1.17, y + 0.49, 0.24, accent)
    # Right output.
    line(slide, 7.35, 4.35, 8.3, 4.35, BLUE, 1.6)
    chevron(slide, 8.06, 4.2, 0.25, 0.3, BLUE)
    line(slide, 8.45, 3.35, 12.35, 3.35, BLUE, 2.0)
    line(slide, 8.45, 5.33, 12.35, 5.33, BLUE, 2.0)
    text(slide, "PERSONALIZED\nLEARNING", 8.75, 3.78, 3.35, 0.7, size=23, color=BLUE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text(slide, "Built in-house  •  Multimodal ingestion  •  Contextual retrieval  •  Long-term learning memory", 0.85, 6.25, 11.7, 0.22, size=8.5, color=FAINT, font=MONO, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "The product experience is powered by a layer we've been building underneath it.\n\nAvenire can take different types of learning material — PDFs, videos, images, web content and notes — and turn them into searchable knowledge.\n\nWhen a student asks something, we don't rely on a single search. We combine different retrieval methods, rerank the results, and use the relevant context to generate the response.\n\nWe've also built the system to recognize when it doesn't have enough evidence rather than confidently making something up.")


def why_now(prs):
    slide = base_slide(prs, 7)
    title_block(slide, "Why now", "AI has commoditized answers.", "The next problem is learning.")
    # Split comparison with center divider.
    line(slide, 6.66, 2.1, 6.66, 5.45, LINE, 1.0)
    text(slide, "BEFORE AI", 1.2, 2.38, 2.6, 0.22, size=10, color=MUTED, bold=True, font=MONO)
    text(slide, "WITH AI", 7.27, 2.38, 2.6, 0.22, size=10, color=BLUE, bold=True, font=MONO)
    left = [("Find information", ORANGE), ("Understand it", TEAL), ("Remember it", PURPLE)]
    for i, (label, accent) in enumerate(left):
        y = 3.0 + i * 0.72
        circle(slide, 1.2, y + 0.07, 0.22, fill=accent)
        text(slide, label, 1.6, y, 3.9, 0.3, size=17, color=TEXT, bold=True)
        if i < 2:
            arrow_down(slide, 1.31, y + 0.31, 0.35, accent)
    text(slide, "Ask", 7.6, 3.1, 3.9, 0.3, size=20, color=TEXT, bold=True)
    arrow_down(slide, 8.0, 3.55, 0.5, BLUE)
    text(slide, "Answer", 7.6, 4.25, 3.9, 0.3, size=20, color=BLUE, bold=True)
    text(slide, "The missing layer is everything after the answer.", M, 6.18, 11.9, 0.42, size=20, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "The timing matters.\n\nAI has made getting an explanation dramatically easier. A student can ask almost anything and get an answer in seconds.\n\nBut getting an answer isn't the same thing as learning.\n\nThe opportunity now is to build the layer that comes after the answer — understanding, practice, retention and personalization.\n\nThat's the space we want Avenire to own.")


def market(prs):
    slide = base_slide(prs, 8)
    title_block(slide, "Market", "Start with the student. Expand into the institution.", "A focused wedge with a natural platform expansion.")
    # Tree.
    node(slide, "AVENIRE", 5.35, 2.0, 2.6, 0.62, fill=BLUE, accent=BLUE, label_size=16)
    line(slide, 6.65, 2.64, 6.65, 3.05, BLUE, 1.8)
    line(slide, 6.65, 3.05, 3.85, 3.55, BLUE, 1.8)
    line(slide, 6.65, 3.05, 9.45, 3.55, TEAL, 1.8)
    chevron(slide, 3.73, 3.42, 0.25, 0.3, BLUE)
    chevron(slide, 9.33, 3.42, 0.25, 0.3, TEAL)
    node(slide, "STUDENTS", 2.45, 3.55, 2.6, 0.62, fill=SURFACE_2, accent=BLUE, label_size=15)
    node(slide, "INSTITUTIONS", 8.25, 3.55, 2.9, 0.62, fill=SURFACE_2, accent=TEAL, label_size=15)
    line(slide, 3.75, 4.2, 3.75, 4.75, BLUE, 1.6)
    line(slide, 9.7, 4.2, 9.7, 4.75, TEAL, 1.6)
    node(slide, "University\nstudents", 2.45, 4.78, 2.6, 0.72, fill=SURFACE_2, accent=BLUE, label_size=12)
    node(slide, "Universities\n& colleges", 8.25, 4.78, 2.9, 0.72, fill=SURFACE_2, accent=TEAL, label_size=12)
    line(slide, 3.75, 5.51, 3.75, 5.85, BLUE, 1.6)
    node(slide, "Self-directed\nlearning", 2.45, 5.9, 2.6, 0.6, fill=SURFACE_3, accent=PURPLE, label_size=11)
    pill(slide, "Initial wedge", 2.45, 6.62, 1.55, color=BLUE, fill=SURFACE_3, size=8)
    pill(slide, "Expansion", 8.25, 6.62, 1.3, color=TEAL, fill=SURFACE_3, size=8)
    add_footer(slide)
    add_notes(slide, "Our initial customer is the student — particularly university students who already learn across many different digital resources.\n\nWe start with a direct-to-student product because it gives us a short feedback loop and lets us understand learning behavior directly.\n\nOnce the product is proven, we see an institutional opportunity as well: universities and colleges can use Avenire as a layer across their existing learning material.\n\nWe're intentionally starting narrow rather than trying to sell to every learner on day one.")


def business_model(prs):
    slide = base_slide(prs, 9)
    title_block(slide, "Business model", "Simple, usage-based pricing.", "The current pricing is a hypothesis we are ready to validate.")
    # Use the actual pricing UI screenshot as product evidence, not as a chart.
    rect(slide, 0.75, 1.95, 8.05, 4.7, fill="080808", line=LINE, radius=False)
    if PRICING.exists():
        slide.shapes.add_picture(str(PRICING), Inches(0.86), Inches(2.08), width=Inches(7.83), height=Inches(4.32))
    line(slide, 9.0, 2.05, 9.0, 6.48, LINE, 1.0)
    text(slide, "CURRENT HYPOTHESIS", 9.6, 2.3, 2.6, 0.22, size=10, color=BLUE, bold=True, font=MONO)
    plans = [("Access", "₹0", "Reduce adoption friction", BLUE), ("Core", "₹450", "Daily study", TEAL), ("Scholar", "₹1,350", "Research-heavy work", ORANGE)]
    for i, (plan, price, desc, accent) in enumerate(plans):
        y = 2.82 + i * 0.95
        circle(slide, 9.62, y + 0.06, 0.18, fill=accent)
        text(slide, plan, 9.95, y, 1.0, 0.2, size=14, color=TEXT, bold=True)
        text(slide, price, 11.05, y - 0.02, 1.1, 0.24, size=16, color=accent, bold=True, align=PP_ALIGN.RIGHT)
        text(slide, desc, 9.95, y + 0.3, 2.1, 0.18, size=10, color=MUTED)
    line(slide, 9.6, 5.75, 12.2, 5.75, PURPLE, 1.0)
    text(slide, "Coming next: Institutional / Enterprise", 9.6, 5.88, 2.6, 0.19, size=9, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "The current business model is straightforward.\n\nWe have a free tier to reduce adoption friction, a ₹450 monthly core plan for regular students, and a ₹1,350 plan for heavier users.\n\nWe're also working toward an institutional tier for universities and organizations.\n\nAt this stage, these are pricing hypotheses rather than proven unit economics — our next goal is validating willingness to pay with real users.")


def traction(prs):
    slide = base_slide(prs, 10)
    title_block(slide, "Traction", "Early validation.", "We are deliberately early — and precise about what is not proven yet.")
    line(slide, 0.85, 2.05, 3.95, 2.05, BLUE, 3.0)
    text(slide, "50+", 0.85, 2.45, 2.5, 0.9, size=52, color=BLUE, bold=True)
    text(slide, "people on our\nwaitlist", 0.9, 3.58, 2.4, 0.6, size=20, color=TEXT, bold=True)
    text(slide, "Where we are today", 4.7, 2.2, 3.4, 0.25, size=11, color=BLUE, bold=True, font=MONO)
    statuses = [("Product", "Working", 1.0, TEAL), ("Waitlist", "50+", 0.76, BLUE), ("Active users", "Pre-launch", 0.08, MUTED), ("Revenue", "Pre-revenue", 0.04, MUTED)]
    for i, (label, value, pct, accent) in enumerate(statuses):
        y = 2.72 + i * 0.68
        text(slide, label, 4.7, y, 1.3, 0.2, size=11, color=TEXT, bold=True)
        add_bar(slide, 6.15, y + 0.02, 3.2, pct, accent, h=0.16)
        text(slide, value, 9.65, y - 0.01, 1.65, 0.22, size=11, color=accent, bold=True, align=PP_ALIGN.RIGHT)
    line(slide, 4.7, 5.62, 12.35, 5.62, LINE, 1.0)
    text(slide, "What we're validating next", 4.95, 5.84, 2.7, 0.2, size=10, color=MUTED, font=MONO)
    text(slide, "Activation  ·  Retention  ·  Willingness to pay  ·  Institutional pilots", 7.5, 5.79, 4.55, 0.28, size=11, color=TEXT, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide)
    add_notes(slide, "We're deliberately early.\n\nWe currently have a working product and a waitlist of around 50 people, but we don't have meaningful active usage or revenue yet.\n\nSo we don't want to present the waitlist as product-market fit.\n\nOur next milestone is to get these first users into the product, measure retention, understand which workflows they actually use, and validate whether they'll pay for it.\n\nThat's exactly the stage where we're looking for incubation.")


def roadmap(prs):
    slide = base_slide(prs, 11)
    title_block(slide, "Roadmap", "From AI study tool to learning system.", "Prove the loop first. Deepen personalization. Then expand into institutions.")
    stages = [
        ("NOW", "Build the learning loop", ["Launch", "Acquire first users", "Measure retention", "Improve core workflows"], BLUE),
        ("NEXT", "Personalize learning", ["Better long-term memory", "Mastery tracking", "Adaptive review", "Personalized study plans"], TEAL),
        ("THEN", "Institutional learning", ["University workspaces", "Institutional knowledge", "Admin / educator tools", "Enterprise deployment"], PURPLE),
    ]
    for i, (phase, head, items, accent) in enumerate(stages):
        x = 0.85 + i * 4.15
        if i > 0:
            line(slide, x - 0.3, 2.05, x - 0.3, 5.98, LINE, 1.0)
        line(slide, x, 2.95, x + 3.05, 2.95, accent, 1.6)
        pill(slide, phase, x + 0.28, 2.33, 0.82, color=accent, fill=SURFACE_3, size=8)
        text(slide, head, x + 0.28, 2.9, 2.9, 0.35, size=17, color=TEXT, bold=True)
        for j, item in enumerate(items):
            y = 3.55 + j * 0.48
            circle(slide, x + 0.3, y + 0.06, 0.13, fill=accent)
            text(slide, item, x + 0.56, y, 2.7, 0.2, size=11, color=MUTED)
        if i < 2:
            line(slide, x + 3.62, 4.0, x + 4.1, 4.0, accent, 1.8)
            chevron(slide, x + 3.98, 3.84, 0.25, 0.3, accent)
    text(slide, "A focused sequence: validation / personalization / infrastructure.", M, 6.35, 11.9, 0.25, size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "Our roadmap follows the same progression.\n\nFirst, we need to prove the core learning loop with individual students.\n\nThen we deepen the personalization — understanding what a student knows, what they've forgotten and what they should work on next.\n\nFinally, we take that system into institutions, where Avenire can sit across the learning material and workflows of an entire university.")


def team(prs):
    slide = base_slide(prs, 12)
    title_block(slide, "Why Avenire", "Built from the student's perspective.", "The founder is already inside the problem — and has built the product from the ground up.")
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(1.15), Inches(2.4), height=Inches(0.72))
    text(slide, "Abhiram", 1.15, 3.55, 2.2, 0.34, size=24, color=TEXT, bold=True)
    text(slide, "Founder", 1.15, 4.05, 2.2, 0.22, size=13, color=BLUE, bold=True, font=MONO)
    text(slide, "BITS Pilani, Hyderabad", 1.15, 4.45, 2.35, 0.24, size=12, color=MUTED)
    line(slide, 4.25, 2.0, 4.25, 5.9, LINE, 1.0)
    bullets = [
        ("Building Avenire from the ground up", "Product, systems and learning workflows"),
        ("Product + engineering", "A working product, not another prototype"),
        ("Student perspective", "Living the fragmented learning workflow every day"),
    ]
    for i, (head, body) in enumerate(bullets):
        y = 2.6 + i * 0.92
        circle(slide, 5.25, y + 0.05, 0.2, fill=[BLUE, TEAL, ORANGE][i])
        text(slide, head, 5.7, y, 5.9, 0.24, size=15, color=TEXT, bold=True)
        text(slide, body, 5.7, y + 0.32, 5.9, 0.22, size=11, color=MUTED)
    text(slide, "What we need now: turn a working product into repeated student behavior.", 4.75, 6.34, 7.6, 0.25, size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "We're building this from inside the problem.\n\nWe're students ourselves, so we experience the fragmented learning workflow we're trying to solve every day.\n\nAt the same time, we've already built the product and the underlying technology ourselves.\n\nWhat we need now isn't another prototype. We need to turn what we've built into something students repeatedly use and eventually pay for.")


def ask(prs):
    slide = base_slide(prs, 13)
    title_block(slide, "The ask", "We're looking for incubation.", "Move Avenire from a working product to a validated company.")
    asks = [
        ("01", "Product validation", "Get our first cohort of students using Avenire.", BLUE),
        ("02", "Mentorship", "Product, education, GTM and business development.", TEAL),
        ("03", "Institutional access", "Pilot with universities and educators.", ORANGE),
        ("04", "Capital", "Infrastructure, development and initial go-to-market.", PURPLE),
    ]
    for i, (num, head, body, accent) in enumerate(asks):
        col = i % 2
        row = i // 2
        x = 0.85 + col * 6.1
        y = 2.05 + row * 1.35
        line(slide, x, y + 1.08, x + 5.6, y + 1.08, LINE, 1.0)
        text(slide, num, x + 0.25, y + 0.24, 0.45, 0.25, size=11, color=accent, bold=True, font=MONO)
        text(slide, head, x + 0.92, y + 0.18, 3.9, 0.26, size=15, color=TEXT, bold=True)
        text(slide, body, x + 0.92, y + 0.54, 4.2, 0.22, size=10.5, color=MUTED)
    rect(slide, 0.85, 5.42, 11.7, 0.88, fill=BLUE, radius=False)
    text(slide, "Our next milestone", 1.15, 5.64, 2.05, 0.2, size=10, color=BG, bold=True, font=MONO)
    text(slide, "50 waitlist", 3.6, 5.58, 1.35, 0.27, size=14, color=BG, bold=True, align=PP_ALIGN.CENTER)
    chevron(slide, 5.1, 5.56, 0.25, 0.3, BG)
    text(slide, "100 active learners", 5.45, 5.58, 1.8, 0.27, size=14, color=BG, bold=True, align=PP_ALIGN.CENTER)
    chevron(slide, 7.4, 5.56, 0.25, 0.3, BG)
    text(slide, "first paying users", 7.75, 5.58, 1.65, 0.27, size=14, color=BG, bold=True, align=PP_ALIGN.CENTER)
    chevron(slide, 9.95, 5.56, 0.25, 0.3, BG)
    text(slide, "institutional pilots", 10.3, 5.58, 1.75, 0.27, size=14, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "This is why we're approaching TBI.\n\nWe have the product foundation, but we're at the point where we need to validate the business.\n\nWe're looking for incubation, mentorship, access to the BITS ecosystem, and capital to help us move from a working product to our first real cohort of users and paying customers.\n\nWe want to use the BITS ecosystem itself as one of our earliest environments for testing and distribution.")


def close(prs):
    slide = base_slide(prs, None)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(M), Inches(1.18), height=Inches(0.52))
    text(slide, "Avenire", M + 0.72, 1.18, 2.1, 0.52, size=27, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    text(slide, "Learning that remembers.", M, 2.25, 11.7, 0.75, size=38, color=TEXT, bold=True)
    text(slide, "The internet gave students access to information.\nAI gave them access to answers.\nWe want to give them something that compounds.", M, 3.35, 7.5, 0.9, size=17, color=MUTED)
    # Small convergence mark.
    for i, (x, y, c) in enumerate([(9.0, 3.0, BLUE), (10.15, 3.78, TEAL), (8.45, 4.55, ORANGE), (10.85, 4.85, PURPLE)]):
        circle(slide, x, y, 0.22, fill=c)
        line(slide, x + 0.11, y + 0.11, 9.9, 4.0, c, 1.1)
    text(slide, "Avenire", M, 6.68, 2.0, 0.22, size=10, color=FAINT, font=MONO)
    text(slide, "abhiram@avenire.ai", W - M - 2.6, 6.68, 2.6, 0.22, size=10, color=FAINT, font=MONO, align=PP_ALIGN.RIGHT)
    add_notes(slide, "The internet gave students access to information.\n\nAI gave them access to answers.\n\nWe want to give them something that compounds:\nan AI that remembers how they learn.\n\nThat's Avenire.\n\nStop.")


def appendix_architecture(prs):
    slide = base_slide(prs, 15)
    title_block(slide, "Appendix A / technical architecture", "From content to learning action.", "The product claim is backed by a concrete, multimodal pipeline.")
    stages = [("Content", BLUE), ("Multimodal\ningestion", TEAL), ("Semantic\nunderstanding", ORANGE), ("Knowledge\nrepresentation", PURPLE), ("Hybrid\nretrieval", PINK), ("Reranking", YELLOW), ("Reasoning", TEAL), ("Learning\naction", BLUE)]
    start_x = 0.55
    box_w = 1.42
    gap = 0.15
    for i, (label, accent) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        y = 3.0 if i % 2 == 0 else 3.8
        rect(slide, x, y, box_w, 0.85, fill=SURFACE_2, line=accent, radius=False)
        text(slide, label, x + 0.12, y + 0.23, box_w - 0.24, 0.38, size=11, color=TEXT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            line(slide, x + box_w, y + 0.43, x + box_w + gap, (3.8 if i % 2 == 0 else 3.0) + 0.43, accent, 1.3)
    text(slide, "Avenire turns scattered learning material and interactions into persistent, searchable, usable knowledge.", 1.0, 5.65, 11.3, 0.5, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Appendix / technical deep dive")
    add_notes(slide, "Appendix: the technical architecture behind the learning intelligence layer.\n\nThe key idea is not a list of vendor technologies. It is the path from raw content to a learning action a student can use.")


def appendix_evaluation(prs):
    slide = base_slide(prs, 16)
    title_block(slide, "Appendix B / retrieval evaluation", "Evidence across difficult query families.", "Small, honest evaluation slices — useful signals, not a claim of solved retrieval.")
    rows = [
        ("Cross-file", 0.607, 0.50, "MRR 60.7%", "evidence 50.0%"),
        ("Multi-hop", 0.592, 0.60, "MRR 59.2%", "evidence 60.0%"),
        ("Formula / chart", 0.514, 0.70, "MRR 51.4%", "evidence 70.0%"),
        ("Timestamp", 0.636, 0.636, "MRR 63.6%", "evidence 63.6%"),
        ("Paraphrase", 0.825, 0.90, "MRR 82.5%", "evidence 90.0%"),
        ("Unanswerable", 0.80, 0.80, "precision 80.0%", "false-positive rate"),
    ]
    text(slide, "FAMILY", 0.85, 2.0, 2.0, 0.2, size=9, color=FAINT, font=MONO)
    text(slide, "MRR@10", 3.45, 2.0, 1.0, 0.2, size=9, color=BLUE, font=MONO)
    text(slide, "ALL REQUIRED EVIDENCE @10", 7.15, 2.0, 2.7, 0.2, size=9, color=TEAL, font=MONO)
    for i, (family, mrr, ev, label1, label2) in enumerate(rows):
        y = 2.35 + i * 0.56
        text(slide, family, 0.85, y, 2.0, 0.22, size=11, color=TEXT, bold=True)
        add_bar(slide, 3.45, y + 0.04, 3.0, mrr, BLUE, h=0.15)
        add_bar(slide, 7.15, y + 0.04, 3.0, ev, TEAL, h=0.15)
        text(slide, label1, 10.35, y - 0.01, 1.35, 0.2, size=9, color=BLUE, font=MONO, align=PP_ALIGN.RIGHT)
        text(slide, label2, 10.35, y + 0.2, 1.35, 0.2, size=8.5, color=TEAL, font=MONO, align=PP_ALIGN.RIGHT)
    rect(slide, 0.85, 5.92, 11.7, 0.56, fill=SURFACE_2, line=LINE, radius=False)
    text(slide, "The hard slices are visible: cross-file and multi-hop are the next engineering frontier.", 1.1, 6.1, 11.2, 0.18, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Appendix / retrieval evaluation")
    add_notes(slide, "Appendix: retrieval evaluation across cross-file, multi-hop, formula/chart, timestamp, paraphrase and unanswerable slices.\n\nThese are engineering signals from small evaluation slices, not a claim that retrieval is solved or that the results generalize to every corpus.")


def appendix_memory(prs):
    slide = base_slide(prs, 17)
    title_block(slide, "Appendix C / learning memory", "The goal is not only to store learning.", "It is to model what a student is likely to remember — and when to bring it back.")
    # Editable curve plot.
    x0, y0, pw, ph = 1.1, 2.05, 8.2, 3.8
    line(slide, x0, y0 + ph, x0 + pw, y0 + ph, LINE, 1.0)
    line(slide, x0, y0, x0, y0 + ph, LINE, 1.0)
    for i in range(1, 5):
        yy = y0 + ph - i * ph / 5
        line(slide, x0, yy, x0 + pw, yy, LINE, 0.6)
    points = []
    for i in range(13):
        x = x0 + (pw / 12) * i
        # Memory stability curve: fast early drop, then a slower plateau.
        val = 0.17 + 0.72 * math.exp(-i / 4.2)
        y = y0 + ph - val * ph
        points.append((x, y))
    for (x1, y1), (x2, y2) in zip(points, points[1:]):
        line(slide, x1, y1, x2, y2, BLUE, 2.4)
    for i, (x, y) in enumerate(points):
        circle(slide, x - 0.055, y - 0.055, 0.11, fill=BLUE if i in (0, 4, 8, 12) else SURFACE_2, line=BLUE)
    text(slide, "memory stability", 0.95, 1.72, 2.0, 0.2, size=10, color=BLUE, font=MONO)
    text(slide, "time", 9.0, 6.03, 0.4, 0.2, size=10, color=FAINT, font=MONO, align=PP_ALIGN.RIGHT)
    # Side insight list: colored rules, not nested cards.
    line(slide, 9.85, 2.1, 12.55, 2.1, TEAL, 1.5)
    text(slide, "Remember", 10.15, 2.38, 2.1, 0.23, size=16, color=TEAL, bold=True)
    text(slide, "What was learned?", 10.15, 2.75, 2.1, 0.2, size=10, color=MUTED)
    line(slide, 9.85, 3.5, 12.55, 3.5, ORANGE, 1.5)
    text(slide, "Predict", 10.15, 3.78, 2.1, 0.23, size=16, color=ORANGE, bold=True)
    text(slide, "What is likely forgotten?", 10.15, 4.15, 2.1, 0.2, size=10, color=MUTED)
    line(slide, 9.85, 4.9, 12.55, 4.9, PURPLE, 1.5)
    text(slide, "Review", 10.15, 5.18, 2.1, 0.23, size=16, color=PURPLE, bold=True)
    text(slide, "When should it return?", 10.15, 5.55, 2.1, 0.2, size=10, color=MUTED)
    add_footer(slide, "Appendix / learning memory")
    add_notes(slide, "Appendix: learning memory.\n\nAvenire is not only storing what a student learned. We are building toward modeling what they are likely to remember, what they are likely to forget, and when a concept should return for review.")


def appendix_pricing(prs):
    slide = base_slide(prs, 18)
    title_block(slide, "Appendix D / pricing", "Current plan structure.", "Pricing is a live hypothesis; the next step is willingness-to-pay validation.")
    rows = [
        ["Plan", "Price / month", "Credits", "Storage", "Learning capabilities"],
        ["Access", "₹0", "330", "2 GB", "AI tutor; interactive concepts; flashcards; misconception detection"],
        ["Core", "₹450", "1,880", "15 GB", "Everything in Access; priority response"],
        ["Scholar", "₹1,350", "6,680", "50 GB", "Everything in Core; mastery tracking; custom study plans; experimental features"],
        ["Coming next", "Institutional / Enterprise", "—", "—", "University workspaces; institutional knowledge; admin and educator tools"],
    ]
    add_table(slide, rows, 0.85, 2.05, 11.7, 2.85, col_widths=[1.55, 1.7, 1.2, 1.2, 6.05], font_size=11)
    line(slide, 0.85, 5.35, 12.55, 5.35, LINE, 1.0)
    text(slide, "Validation question", 1.15, 5.57, 1.8, 0.2, size=10, color=BLUE, bold=True, font=MONO)
    text(slide, "Which workflow creates enough repeated value for students to pay — and for institutions to pilot?", 3.2, 5.52, 8.9, 0.3, size=15, color=TEXT, bold=True)
    add_footer(slide, "Appendix / pricing")
    add_notes(slide, "Appendix: the full current pricing structure.\n\nThe important caveat is that these are pricing hypotheses, not proven unit economics. The next step is to put the plans in front of real users and learn which workflows create enough repeated value to pay for.")


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "Avenire — Learning that remembers"
    prs.core_properties.subject = "TBI incubation pitch deck"
    prs.core_properties.author = "Abhiram"
    prs.core_properties.company = "Avenire"
    prs.core_properties.comments = "Editable PowerPoint deck. Charts are native shapes and the pricing appendix uses a real PowerPoint table."

    cover(prs)
    problem(prs)
    insight(prs)
    product(prs)
    loop(prs)
    technology(prs)
    why_now(prs)
    market(prs)
    business_model(prs)
    traction(prs)
    roadmap(prs)
    team(prs)
    ask(prs)
    close(prs)
    appendix_architecture(prs)
    appendix_evaluation(prs)
    appendix_memory(prs)
    appendix_pricing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
