from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


INPUT = Path(__file__).with_name("avenire_pitch_deck.pptx")
OUTPUT = INPUT
EMU = 914400


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def add_text(slide, text, x, y, w, h, *, font="Inter", size=11, color="E4E4E4",
             bold=False, align=PP_ALIGN.LEFT, margin=0, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(int(x * EMU), int(y * EMU), int(w * EMU), int(h * EMU))
    box.name = "Redesign text"
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = int(margin * EMU)
    tf.margin_top = tf.margin_bottom = int(margin * EMU)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = 0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = int(size * 12700)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rule(slide, x, y, w, color, thickness=1.0):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        int(x * EMU), int(y * EMU), int(w * EMU), int(thickness / 72 * EMU),
    )
    line.name = "Redesign rule"
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(color)
    line.line.fill.background()
    return line


def add_flat_strip(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        int(x * EMU), int(y * EMU), int(w * EMU), int(h * EMU),
    )
    shape.name = "Redesign conclusion strip"
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = int(0.75 * 12700)
    else:
        shape.line.fill.background()
    return shape


prs = Presentation(INPUT)
slide = prs.slides[3]

# Preserve the shared header/footer and remove only the old comparison treatment.
for shape in list(slide.shapes):
    if shape.name not in {"Rectangle 1", "Picture 2", "TextBox 3", "TextBox 4",
                          "TextBox 5", "TextBox 6", "TextBox 7", "Connector 28",
                          "TextBox 29"}:
        remove_shape(shape)

# Three columns share a single visual system. A slim accent rule replaces each
# outer card and the repeated rounded callout.
columns = [
    (1.00, 3.55, "78C4D5", "ChatGPT", "Answers anything, instantly.", "No learner model"),
    (4.95, 3.55, "A99AD1", "Notion AI", "Everything in one place.", "No spaced repetition"),
    (8.90, 3.55, "F0A873", "Anki", "Flashcards built for memory.", "No misconception diagnosis"),
]

for x, width, accent, name, description, gap in columns:
    add_rule(slide, x, 2.60, width, accent, thickness=2.0)
    add_text(slide, name, x, 2.86, width, 0.34, size=18, bold=True)
    add_text(slide, description, x, 3.36, width, 0.28, size=11, color="9C9C9C")
    add_rule(slide, x, 4.02, width, "444444", thickness=0.75)
    add_text(slide, "Missing", x, 4.30, width, 0.18, font="Cascadia Mono", size=9,
             color="FC6B83", bold=True)
    add_text(slide, gap, x, 4.53, width, 0.32, size=11, color="E4E4E4", bold=True)

# One restrained conclusion strip anchors the slide and keeps the takeaway
# visually distinct without adding another rounded card.
add_flat_strip(slide, 0.85, 5.92, 0.06, 0.88, "ABC4FF")
add_text(
    slide,
    "Avenire is the only one that gets smarter about you — not just the material.",
    1.15, 6.15, 11.35, 0.40, size=17, color="E4E4E4", bold=True,
    valign=MSO_ANCHOR.MIDDLE,
)

prs.save(OUTPUT)
print(f"saved {OUTPUT}")
