from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from create_avenire_pitch import (
    BG,
    BLUE,
    LINE,
    MUTED,
    ORANGE,
    PINK,
    PURPLE,
    TEAL,
    TEXT,
    W,
    H,
    add_footer,
    add_notes,
    brand,
    circle,
    line,
    rect,
    text,
    title_block,
)


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "pitch_deck" / "avenire_pitch_deck.pptx"


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def replace_text(slide, old: str, new: str):
    for shape in slide.shapes:
        if getattr(shape, "text", "") == old and shape.has_text_frame:
            paragraph = shape.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].text = new
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = new
            return True
    return False


def reposition_source_logos(prs):
    """Keep the user's added logo objects, but make them read as part of the nodes."""
    slide = prs.slides[1]
    pictures = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.left > Inches(5)
    ]
    # The YouTube logo is the upper picture; the ChatGPT logo is the right picture.
    for picture in pictures:
        if picture.top < Inches(2.3):
            picture.left = Inches(6.14)
            picture.top = Inches(1.82)
            picture.width = Inches(0.25)
            picture.height = Inches(0.22)
        else:
            picture.left = Inches(10.27)
            picture.top = Inches(2.69)
            picture.width = Inches(0.24)
            picture.height = Inches(0.24)
    for shape in slide.shapes:
        if getattr(shape, "text", "") == "YouTube":
            shape.left = Inches(6.49)
            shape.top = Inches(1.82)
            shape.width = Inches(0.92)
        elif getattr(shape, "text", "") == "ChatGPT":
            shape.left = Inches(10.58)
            shape.top = Inches(2.72)
            shape.width = Inches(0.98)


def rebuild_traction_slide(prs):
    """Replace the ambiguous bars with a clear, editable status path."""
    slide = prs.slides[9]
    for shape in list(slide.shapes):
        remove_shape(shape)
    rect(slide, 0, 0, W, H, fill=BG)
    brand(slide, 10)
    title_block(
        slide,
        "Traction",
        "Early validation.",
        "We are deliberately early — and precise about what is not proven yet.",
    )

    line(slide, 0.85, 2.12, 3.15, 2.12, BLUE, 3.0)
    text(slide, "50+", 0.85, 2.5, 2.5, 0.9, size=52, color=BLUE, bold=True)
    text(slide, "people on our waitlist", 0.9, 3.62, 2.9, 0.3, size=17, color=TEXT, bold=True)
    text(slide, "A working product with early interest — before active usage or revenue.", 0.9, 4.25, 2.8, 0.65, size=11, color=MUTED)

    text(slide, "Where we are today", 4.55, 2.18, 3.4, 0.25, size=11, color=BLUE, bold=True, font="Cascadia Mono")
    line(slide, 4.75, 3.28, 11.95, 3.28, LINE, 1.6)
    statuses = [
        ("Product", "Working", TEAL),
        ("Waitlist", "50+", BLUE),
        ("Active users", "Pre-launch", MUTED),
        ("Revenue", "Pre-revenue", MUTED),
    ]
    xs = [4.95, 7.05, 9.15, 11.25]
    for i, ((label, value, accent), x) in enumerate(zip(statuses, xs)):
        circle(slide, x, 3.08, 0.4, fill=accent, line=accent if accent != MUTED else LINE)
        text(slide, str(i + 1), x, 3.17, 0.4, 0.15, size=8, color=BG, bold=True, font="Cascadia Mono", align=2)
        text(slide, label, x - 0.45, 3.78, 1.3, 0.2, size=11, color=TEXT, bold=True, align=2)
        text(slide, value, x - 0.6, 4.12, 1.6, 0.2, size=10, color=accent, bold=True, align=2)

    line(slide, 4.55, 5.45, 12.05, 5.45, LINE, 1.0)
    text(slide, "What we're validating next", 4.55, 5.7, 2.7, 0.2, size=10, color=MUTED, font="Cascadia Mono")
    text(slide, "Activation  ·  Retention  ·  Willingness to pay  ·  Institutional pilots", 7.25, 5.65, 4.8, 0.28, size=11, color=TEXT, bold=True, align=2)
    add_footer(slide)
    add_notes(
        slide,
        "We're deliberately early.\n\nWe currently have a working product and a waitlist of around 50 people, but we don't have meaningful active usage or revenue yet.\n\nSo we don't want to present the waitlist as product-market fit.\n\nThis slide is intentionally a status path, not a performance chart: product working, waitlist 50+, active users pre-launch, and revenue pre-revenue.\n\nOur next milestone is to get these first users into the product, measure retention, understand which workflows they actually use, and validate whether they'll pay for it.",
    )


def enrich_product_slide(prs):
    slide = prs.slides[3]
    # Replace only the open right-side annotation area; keep the user's UI work intact.
    for shape in list(slide.shapes):
        if shape.left >= Inches(9.0) and Inches(1.8) < shape.top < Inches(6.8):
            remove_shape(shape)
    line(slide, 9.1, 2.18, 9.1, 6.42, LINE, 1.0)
    features = [
        ("Visualize", "Turn abstract concepts into interactive models.", BLUE),
        ("Review", "Use spaced repetition to bring back what matters.", TEAL),
        ("Diagnose", "Spot misconception patterns while the engine is being built.", ORANGE),
    ]
    for i, (head, body, accent) in enumerate(features):
        y = 2.35 + i * 1.2
        line(slide, 9.55, y, 12.25, y, accent, 1.5)
        text(slide, head, 9.55, y + 0.18, 2.4, 0.25, size=16, color=accent, bold=True)
        text(slide, body, 9.55, y + 0.53, 2.65, 0.42, size=10.5, color=MUTED)
    add_notes(
        slide,
        "This is Avenire.\n\nWe call the experience Conversational Learning.\n\nA student can ask a question, explore an explanation, interact with a visualization, create notes, generate practice material and eventually review that concept again.\n\nWe have also been working out three learning actions that make the experience compound: visualizations for abstract concepts, spaced repetition for bringing concepts back at the right time, and a misconception engine that is currently being built to diagnose why an answer is wrong rather than only marking it wrong.\n\nThese aren't separate workflows. They're connected parts of the same learning process.",
    )


def enrich_learning_loop(prs):
    slide = prs.slides[4]
    replace_text(slide, "Memory", "Memory + review")
    replace_text(slide, "One continuous learning loop.", "One continuous learning loop — with review that compounds.")
    text(slide, "Misconception engine", 0.9, 5.0, 2.6, 0.22, size=13, color=ORANGE, bold=True)
    text(slide, "Diagnose what is wrong, not only what was missed.", 0.9, 5.32, 4.25, 0.22, size=10.5, color=MUTED)
    text(slide, "Spaced repetition", 8.3, 5.0, 2.6, 0.22, size=13, color=TEAL, bold=True)
    text(slide, "Bring the right concept back when memory needs it.", 8.3, 5.32, 4.0, 0.22, size=10.5, color=MUTED)
    add_notes(
        slide,
        "This is the fundamental difference between Avenire and a generic AI assistant.\n\nWith a normal chatbot, I ask a question, get an answer, and move on.\n\nWith Avenire, that interaction can become part of my learning system. A conversation can become knowledge. Knowledge can become practice. Practice can become memory.\n\nThe next layer is what makes the loop useful over time: spaced repetition brings the right concept back when it is likely to be forgotten, while the misconception engine we are building will diagnose patterns in wrong answers.\n\nAnd visualizations make abstract ideas something the student can manipulate, not only read about.\n\nThat is one continuous learning loop — with review that compounds.",
    )


def enrich_roadmap(prs):
    slide = prs.slides[10]
    replace_text(slide, "Adaptive review", "Spaced repetition")
    text(slide, "Misconception engine / in build", 1.1, 5.28, 2.8, 0.22, size=10, color=ORANGE, bold=True)
    replace_text(
        slide,
        "A focused sequence: validation / personalization / infrastructure.",
        "Now building: misconception diagnosis  ·  Next: spaced repetition + mastery tracking.",
    )
    add_notes(
        slide,
        "Our roadmap follows the same progression.\n\nFirst, we need to prove the core learning loop with individual students, including the misconception engine that is currently being built.\n\nThen we deepen personalization: spaced repetition, mastery tracking, and understanding what a student knows, what they've forgotten and what they should work on next.\n\nFinally, we take that system into institutions, where Avenire can sit across the learning material and workflows of an entire university.",
    )


def enrich_memory_appendix(prs):
    slide = prs.slides[16]
    replace_text(slide, "Remember", "Spaced repetition")
    replace_text(slide, "What was learned?", "Schedule the next useful review.")
    replace_text(slide, "Predict", "Misconception engine")
    replace_text(slide, "What is likely forgotten?", "Diagnose patterns in wrong answers.")
    replace_text(slide, "Review", "Visualizations")
    replace_text(slide, "When should it return?", "Make abstract ideas manipulable.")
    add_notes(
        slide,
        "Appendix: learning memory and the next layer of the product.\n\nAvenire is not only storing what a student learned. We are building toward modeling what they are likely to remember, what they are likely to forget, and when a concept should return for review.\n\nSpaced repetition is the mechanism for timing that review. The misconception engine is being built to identify patterns in wrong answers. And interactive visualizations help students form a mental model before the review loop begins.",
    )


def normalize_legacy_background(prs):
    # The edited source deck carried a slightly oversized slide-6 background.
    # Keep the visual treatment, but bring the shape back inside the canvas.
    slide = prs.slides[5]
    if slide.shapes:
        background = slide.shapes[0]
        background.left = Inches(0)
        background.top = Inches(0)
        background.width = Inches(W)
        background.height = Inches(H)


def main():
    prs = Presentation(DECK)
    reposition_source_logos(prs)
    rebuild_traction_slide(prs)
    enrich_product_slide(prs)
    enrich_learning_loop(prs)
    enrich_roadmap(prs)
    enrich_memory_appendix(prs)
    normalize_legacy_background(prs)
    temp = DECK.with_name(".avenire_pitch_deck.updated.pptx")
    prs.save(temp)
    temp.replace(DECK)
    print(DECK)


if __name__ == "__main__":
    main()
