from pathlib import Path

from pptx import Presentation


deck = Path(__file__).resolve().parent / "avenire_pitch_deck.pptx"
prs = Presentation(deck)

# Current deck has the live-demo slide at position 5. Move it after the close
# so the original main-pitch slide numbers remain stable for the presenter.
demo = prs.slides[4]
slide_id = next(item for item in prs.slides._sldIdLst if item.id == demo.slide_id)
prs.slides._sldIdLst.remove(slide_id)
prs.slides._sldIdLst.insert(14, slide_id)

for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        value = getattr(shape, "text", "")
        if value.isdigit() and len(value) == 2 and shape.left > 11.5 * 914400 and shape.top < 0.8 * 914400:
            shape.text_frame.paragraphs[0].runs[0].text = f"{i:02d}"

temp = deck.with_name(".avenire_pitch_deck.repositioned.pptx")
prs.save(temp)
temp.replace(deck)
print(deck)
