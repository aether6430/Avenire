from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from create_avenire_pitch import (
    BG,
    BLUE,
    FAINT,
    LINE,
    MUTED,
    ORANGE,
    PINK,
    PURPLE,
    TEAL,
    TEXT,
    W,
    H,
    add_footer,
    add_notes,
    brand,
    chevron,
    circle,
    line,
    rect,
    text,
    title_block,
)
from update_current_pitch import (
    enrich_learning_loop,
    enrich_memory_appendix,
    enrich_product_slide,
    enrich_roadmap,
    normalize_legacy_background,
    reposition_source_logos,
    rebuild_traction_slide,
    remove_shape,
)


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "pitch_deck" / "avenire_pitch_deck.pptx"


def clear_slide(slide):
    for shape in list(slide.shapes):
        remove_shape(shape)


def set_notes(slide, notes):
    """Write notes even when a user-edited deck has an empty notes master."""
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is not None:
        notes_slide.notes_text_frame.text = notes
        return
    sp_tree = notes_slide._element.xpath(".//p:spTree")[0]
    sp = OxmlElement("p:sp")
    nv_sp_pr = OxmlElement("p:nvSpPr")
    c_nv_pr = OxmlElement("p:cNvPr")
    c_nv_pr.set("id", "10")
    c_nv_pr.set("name", "Notes Placeholder")
    c_nv_sp_pr = OxmlElement("p:cNvSpPr")
    nv_pr = OxmlElement("p:nvPr")
    ph = OxmlElement("p:ph")
    ph.set("type", "body")
    ph.set("idx", "3")
    nv_pr.append(ph)
    nv_sp_pr.extend([c_nv_pr, c_nv_sp_pr, nv_pr])
    sp.append(nv_sp_pr)
    sp.append(OxmlElement("p:spPr"))
    tx_body = OxmlElement("p:txBody")
    tx_body.append(OxmlElement("a:bodyPr"))
    tx_body.append(OxmlElement("a:lstStyle"))
    for paragraph_text in notes.split("\n"):
        paragraph = OxmlElement("a:p")
        run = OxmlElement("a:r")
        run_props = OxmlElement("a:rPr")
        run.append(run_props)
        value = OxmlElement("a:t")
        value.text = paragraph_text
        run.append(value)
        paragraph.append(run)
        paragraph.append(OxmlElement("a:endParaRPr"))
        tx_body.append(paragraph)
    sp.append(tx_body)
    sp_tree.append(sp)


def replace_text(slide, old: str, new: str):
    for shape in slide.shapes:
        if getattr(shape, "text", "") == old and shape.has_text_frame:
            paragraph = shape.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].text = new
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = new
            return True
    return False


def rebuild_recall_problem(prs):
    slide = prs.slides[2]
    clear_slide(slide)
    rect(slide, 0, 0, W, H, fill=BG)
    brand(slide, 3)
    title_block(
        slide,
        "The missing piece",
        "I read this somewhere. But I forgot.",
        "Students remember the feeling of understanding — not where the idea lived or how it connected.",
    )
    text(slide, "The next search starts\nwithout the context\nthat made the idea useful.", 0.9, 2.35, 4.6, 1.45, size=28, color=TEXT, bold=True)
    line(slide, 0.9, 4.25, 4.95, 4.25, ORANGE, 2.0)
    text(slide, "Avenire turns scattered learning material into\nsearchable, connected evidence.", 0.9, 4.5, 4.7, 0.6, size=15, color=MUTED)

    text(slide, "Retrieval that brings the idea back", 6.1, 2.35, 5.7, 0.28, size=16, color=BLUE, bold=True)
    labels = [
        ("Scattered\nmaterial", BLUE),
        ("Retrieve\nevidence", TEAL),
        ("Reconnect\ncontext", ORANGE),
        ("Practice +\nremember", PURPLE),
    ]
    xs = [6.05, 7.7, 9.35, 11.0]
    for i, ((label, accent), x) in enumerate(zip(labels, xs)):
        circle(slide, x, 3.32, 0.62, fill=accent)
        text(slide, str(i + 1), x, 3.49, 0.62, 0.18, size=10, color=BG, bold=True, font="Cascadia Mono", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        text(slide, label, x - 0.28, 4.14, 1.18, 0.48, size=10.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(labels) - 1:
            line(slide, x + 0.72, 3.63, xs[i + 1] - 0.15, 3.63, accent, 1.5)
            chevron(slide, xs[i + 1] - 0.23, 3.5, 0.2, 0.27, accent)
    line(slide, 6.05, 5.25, 12.0, 5.25, LINE, 1.0)
    text(slide, "Semantic  ·  lexical  ·  exact-match  ·  reranking  ·  evidence threshold", 6.05, 5.5, 5.95, 0.2, size=9.5, color=MUTED, font="Cascadia Mono", align=PP_ALIGN.CENTER)
    add_footer(slide)
    set_notes(
        slide,
        "One of the most common learning failures is: I read this somewhere, but I forgot.\n\nThe student does not only lose the answer. They lose the source, the surrounding context and the path back to the idea.\n\nThat is why retrieval is central to Avenire. We combine multiple retrieval signals across PDFs, videos, notes, chats and other material, then rerank the evidence and use an evidence threshold so the system can abstain when it cannot support an answer.\n\nThe goal is not just to answer again. It is to reconnect the student with the context so the next interaction can become practice and memory.",
    )


def rebuild_business_model(prs):
    slide = prs.slides[8]
    clear_slide(slide)
    rect(slide, 0, 0, W, H, fill=BG)
    brand(slide, 9)
    title_block(
        slide,
        "Business model",
        "Simple, usage-based pricing.",
        "A current hypothesis designed to reduce adoption friction and learn what students will pay for.",
    )
    cols = [
        (0.85, "Access", "₹0", "For students getting started.", ["330 Apollo credits", "2 GB storage", "AI tutor", "Interactive concepts", "Flashcards"], BLUE),
        (4.55, "Core", "₹450", "For daily study.", ["1,880 Apollo credits", "15 GB storage", "Priority response", "Everything in Access"], TEAL),
        (8.25, "Scholar", "₹1,350", "For research-heavy work.", ["6,680 Apollo credits", "50 GB storage", "Mastery tracking", "Custom study plans", "Experimental features"], ORANGE),
    ]
    for i, (x, name, price, desc, features, accent) in enumerate(cols):
        if i:
            line(slide, x - 0.28, 2.05, x - 0.28, 5.85, LINE, 1.0)
        line(slide, x, 2.35, x + 2.9, 2.35, accent, 2.0)
        text(slide, name, x, 2.58, 2.4, 0.3, size=20, color=TEXT, bold=True)
        text(slide, price, x, 3.08, 2.4, 0.45, size=28, color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)
        text(slide, desc, x, 3.65, 2.9, 0.3, size=11, color=MUTED)
        for j, feature in enumerate(features):
            y = 4.35 + j * 0.36
            circle(slide, x, y + 0.04, 0.11, fill=accent)
            text(slide, feature, x + 0.25, y, 2.65, 0.2, size=10.5, color=TEXT, valign=MSO_ANCHOR.MIDDLE)
    line(slide, 0.85, 6.15, 12.45, 6.15, PURPLE, 1.2)
    text(slide, "Coming next", 0.85, 6.35, 1.25, 0.2, size=10, color=PURPLE, bold=True, font="Cascadia Mono")
    text(slide, "Institutional / Enterprise — university workspaces, educator tools and deployment across institutional knowledge.", 2.25, 6.32, 10.1, 0.28, size=11, color=TEXT)
    add_footer(slide)
    set_notes(
        slide,
        "The current business model is straightforward.\n\nWe have a free tier to reduce adoption friction, a ₹450 monthly core plan for regular students, and a ₹1,350 plan for heavier users.\n\nWe are also working toward an institutional tier for universities and organizations.\n\nAt this stage, these are pricing hypotheses rather than proven unit economics. The next goal is validating willingness to pay with real users.",
    )


def enrich_retrieval_slide(prs):
    slide = prs.slides[5]
    replace_text(
        slide,
        "Built in-house  •  Multimodal ingestion  •  Contextual retrieval  •  Long-term learning memory",
        "Retrieval stack  •  semantic + lexical + exact-match  •  reranking  •  abstention",
    )
    add_notes(
        slide,
        "The product experience is powered by a learning intelligence layer we have been building underneath it.\n\nAvenire can take PDFs, videos, images, web content and notes and turn them into searchable knowledge.\n\nThe retrieval work is deliberately multi-signal: semantic search for meaning, lexical and exact-match paths for names and formulas, fusion and reranking for evidence quality, and an abstention threshold when the evidence is not sufficient.\n\nOur research work is about making that evidence path measurable across difficult slices such as cross-file, multi-hop, formula/chart, timestamp, paraphrase and unanswerable queries.",
    )


def add_live_demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, fill=BG)
    brand(slide, 5)
    title_block(slide, "Live demo", "The visualization engine.", "Avenire turns abstract ideas into things a student can manipulate.")
    rect(slide, 0.9, 2.15, 11.55, 3.55, fill=BG, line=LINE, radius=False)
    text(slide, "LIVE DEMO", 1.1, 2.42, 1.3, 0.2, size=10, color=BLUE, bold=True, font="Cascadia Mono")
    text(slide, "Show the engine here.", 1.1, 3.2, 11.1, 0.58, size=32, color=TEXT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text(slide, "Double pendulum  ·  DNA double helix  ·  interactive concept models", 1.1, 4.25, 11.1, 0.25, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    line(slide, 1.1, 5.1, 12.2, 5.1, TEAL, 1.6)
    text(slide, "Not a screenshot: open the real product and let the room see the model respond.", 1.1, 5.28, 11.1, 0.22, size=10.5, color=TEAL, align=PP_ALIGN.CENTER)
    add_footer(slide)
    set_notes(
        slide,
        "This is the point where I would switch from the deck to the live product.\n\nShow one or two interactive concepts — for example a double pendulum or a DNA double helix — and demonstrate that the student can manipulate the model, not just look at a static image.\n\nThe point is not visual polish by itself. The visualization is a learning action: it gives the student a mental model they can question, connect to notes and return to during review.",
    )
    return slide


def add_research_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, fill=BG)
    brand(slide, len(prs.slides))
    title_block(
        slide,
        "Appendix E / research",
        "Separating topic overlap from misconception.",
        "A proposed direction for the misconception engine — research in progress, not a shipped claim.",
    )
    # Editable vector sketch.
    ox, oy = 1.25, 5.65
    line(slide, ox, oy, 6.65, oy, LINE, 1.0)
    line(slide, ox, oy, ox, 2.65, LINE, 1.0)
    text(slide, "topic familiarity", 4.85, 5.82, 1.7, 0.2, size=9.5, color=FAINT, font="Cascadia Mono", align=PP_ALIGN.RIGHT)
    text(slide, "belief signal", 0.8, 2.55, 1.4, 0.2, size=9.5, color=FAINT, font="Cascadia Mono")
    line(slide, ox, oy, 5.0, 4.7, TEAL, 1.6)
    line(slide, ox, oy, 4.55, 3.35, ORANGE, 2.0)
    line(slide, ox, oy, 6.0, 3.0, PINK, 2.0)
    circle(slide, 4.87, 4.57, 0.18, fill=TEAL)
    circle(slide, 4.42, 3.22, 0.18, fill=ORANGE)
    circle(slide, 5.87, 2.87, 0.18, fill=PINK)
    text(slide, "topic-only question", 5.1, 4.48, 1.45, 0.2, size=10.5, color=TEAL)
    text(slide, "known wrong belief", 4.1, 3.0, 1.55, 0.2, size=10.5, color=ORANGE, align=PP_ALIGN.RIGHT)
    text(slide, "novel confusion", 5.98, 2.65, 1.3, 0.2, size=10.5, color=PINK)
    text(slide, "δmis = full_emb − topic_emb", 1.25, 6.15, 3.2, 0.2, size=11, color=ORANGE, font="Cascadia Mono")

    line(slide, 8.05, 2.2, 8.05, 6.25, LINE, 1.0)
    sections = [
        ("Current system", "Cosine prefilter over full records, then an online LLM classifier on borderline hits.", BLUE),
        ("Research direction", "Store topic_emb and full_emb; cancel topic overlap and score the misconception delta.", ORANGE),
        ("Novelty signal", "Residual magnitude surfaces confusion that no known misconception explains.", PINK),
        ("Why it matters", "Reject topic bleed, reduce online classifier work, and send novel cases to post-hoc review instead of hallucinating mid-turn.", TEAL),
    ]
    for i, (head, body, accent) in enumerate(sections):
        y = 2.35 + i * 0.98
        line(slide, 8.45, y, 11.95, y, accent, 1.5)
        text(slide, head, 8.45, y + 0.17, 3.0, 0.23, size=14, color=accent, bold=True)
        text(slide, body, 8.45, y + 0.48, 3.65, 0.4, size=10.5, color=MUTED)
    add_footer(slide, "Appendix / misconception research")
    set_notes(
        slide,
        "This is a research direction we are actively thinking through for the misconception engine.\n\nThe current approach embeds a full misconception record, uses cosine similarity to find candidates, and then calls an LLM classifier to decide whether the learner is actually expressing that misconception. The risk is topic bleed: a normal question about Newton's third law can look similar to a record about a specific wrong belief.\n\nThe proposed direction stores two anchors: a topic embedding and a full misconception embedding. Their difference gives a misconception delta direction. We can project away the topic component, score alignment with the delta, and use the residual magnitude as a novelty signal.\n\nKnown matches can be handled online. Novel or ambiguous cases can be flagged for post-hoc review, where the system can create a new record without hallucinating an intervention in the middle of the tutoring turn.\n\nThis is proposal-stage research, with thresholds and stability requirements still to be validated.",
    )
    return slide


def center_flow_text(prs):
    flow_texts = {
        "PDF", "YouTube", "ChatGPT", "Notes", "Flashcards", "student",
        "ASK", "UNDERSTAND", "PRACTICE", "REMEMBER",
        "UNDERSTAND", "CONNECT", "RETRIEVE", "REASON",
        "AVENIRE", "STUDENTS", "INSTITUTIONS", "University\nstudents", "Universities\n& colleges", "Self-directed\nlearning",
        "Content", "Multimodal\ningestion", "Semantic\nunderstanding", "Knowledge\nrepresentation", "Hybrid\nretrieval", "Reranking", "Reasoning", "Learning\naction",
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "text", "") in flow_texts:
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def move_slide(prs, slide, index):
    slide_id = next(item for item in prs.slides._sldIdLst if item.id == slide.slide_id)
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(index, slide_id)


def renumber(prs):
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            value = getattr(shape, "text", "")
            if re.fullmatch(r"\d{2}", value) and shape.left > Inches(11.5) and shape.top < Inches(0.8):
                shape.text_frame.paragraphs[0].runs[0].text = f"{i:02d}"


def main():
    prs = Presentation(DECK)

    # Apply the requested changes to the existing 18-slide deck before inserting new slides.
    reposition_source_logos(prs)
    rebuild_recall_problem(prs)
    enrich_product_slide(prs)
    enrich_learning_loop(prs)
    enrich_retrieval_slide(prs)
    rebuild_business_model(prs)
    rebuild_traction_slide(prs)
    enrich_roadmap(prs)
    enrich_memory_appendix(prs)
    normalize_legacy_background(prs)
    center_flow_text(prs)

    demo = add_live_demo_slide(prs)
    move_slide(prs, demo, 14)  # after the close, keeping the main narrative numbering stable
    add_research_slide(prs)
    center_flow_text(prs)
    renumber(prs)

    temp = DECK.with_name(".avenire_pitch_deck.v3.pptx")
    prs.save(temp)
    temp.replace(DECK)
    print(DECK)


if __name__ == "__main__":
    main()
