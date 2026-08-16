from __future__ import annotations

import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_avenire_pitch import (
    BG,
    BLUE,
    FAINT,
    LINE,
    MUTED,
    ORANGE,
    PURPLE,
    RED,
    SURFACE_2,
    SURFACE_3,
    TEAL,
    TEXT,
    MONO,
    add_footer,
    base_slide,
    line,
    rect,
    text,
    title_block,
)
from update_current_pitch import remove_shape
from update_pitch_v3 import move_slide, set_notes


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "pitch_deck" / "avenire_pitch_deck.pptx"


def find_slide(prs, kicker):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "text", "") == kicker:
                return slide
    raise ValueError(f"No slide with kicker {kicker!r}")


def find_group_containing(slide, needle):
    def recurse(shape):
        if shape.shape_type == 6:
            if any(getattr(c, "text", "") == needle for c in shape.shapes):
                return shape
            for child in shape.shapes:
                found = recurse(child)
                if found is not None:
                    return found
        return None

    for shape in slide.shapes:
        found = recurse(shape)
        if found is not None:
            return found
    raise ValueError(f"No group containing {needle!r}")


def collect_text_shapes(shape, needle, out):
    if getattr(shape, "text", "") == needle:
        out.append(shape)
    if shape.shape_type == 6:
        for child in shape.shapes:
            collect_text_shapes(child, needle, out)


def replace_text(slide, old, new):
    for shape in slide.shapes:
        if getattr(shape, "text", "") == old and shape.has_text_frame:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = new
                for run in p.runs[1:]:
                    run.text = ""
            else:
                p.text = new
            return True
    return False


def renumber(prs):
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            value = getattr(shape, "text", "")
            if re.fullmatch(r"\d{2}", value) and shape.left > Inches(11.5) and shape.top < Inches(0.8):
                shape.text_frame.paragraphs[0].runs[0].text = f"{i:02d}"


def cover_hook(prs):
    slide = prs.slides[0]
    text(slide, "Working product  ·  50+ waitlist", 0.65, 5.95, 6.4, 0.35, size=16, color=BLUE, bold=True)
    set_notes(slide, add_cover_notes(slide))


def add_cover_notes(slide):
    existing = slide.notes_slide.notes_text_frame.text
    hook = "\n\nNow: a working product and 50+ on the waitlist. That is the hook this pitch is built to validate."
    return existing + hook


def rebuild_missing_piece_right(prs):
    slide = find_slide(prs, "THE MISSING PIECE")
    group = find_group_containing(slide, "Retrieval that brings the idea back")
    remove_shape(group)

    text(slide, "Retrieval that brings the idea back", 6.1, 2.6, 6.1, 0.3, size=16, color=BLUE, bold=True)
    rect(slide, 6.1, 3.15, 6.1, 2.1, fill=SURFACE_2, line=LINE, radius=True)
    text(
        slide,
        "You highlighted a paragraph on backprop 3 weeks ago — Avenire resurfaces it when you ask about vanishing gradients.",
        6.35,
        3.4,
        5.6,
        1.6,
        size=15,
        color=TEXT,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    text(
        slide,
        "Semantic  ·  lexical  ·  exact-match  ·  reranking  ·  evidence threshold",
        6.1,
        5.5,
        6.1,
        0.2,
        size=9.5,
        color=MUTED,
        font=MONO,
        align=PP_ALIGN.CENTER,
    )
    set_notes(
        slide,
        "One of the most common learning failures: I read this somewhere, but I forgot.\n\nThe student does not only lose the answer. They lose the source, the surrounding context and the path back to the idea.\n\nA concrete example carries more than a generic flow: you highlighted a paragraph on backprop 3 weeks ago, and Avenire resurfaces it when you ask about vanishing gradients.\n\nThat reconnecting of context — across PDFs, videos, notes and chats — is what retrieval means here, backed by an evidence threshold so the system abstains when it cannot support an answer.",
    )


def technology_proof(prs):
    slide = find_slide(prs, "THE TECHNOLOGY")
    text(
        slide,
        "82.5% MRR on paraphrase queries, 60%+ on cross-file retrieval",
        0.85,
        5.9,
        11.7,
        0.35,
        size=14,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    set_notes(
        slide,
        "The product experience is powered by a learning intelligence layer built underneath it.\n\nAvenire takes PDFs, videos, images, web content and notes and turns them into searchable knowledge.\n\nThe retrieval work is multi-signal: semantic search for meaning, lexical and exact-match paths for names and formulas, fusion and reranking for evidence quality, and an abstention threshold when evidence is not sufficient.\n\nMeasured on our own eval slices: 82.5% MRR on paraphrase queries, 60%+ on cross-file retrieval. Full breakdown in Appendix B. These are useful engineering signals, not a claim that retrieval is solved.",
    )


def competitive_landscape(prs):
    slide = base_slide(prs, 1)
    title_block(
        slide,
        "Competitive landscape",
        "Everyone has a tool. Nobody has the learner.",
        "Three tools students already use — each solves one slice, none models the student.",
    )
    cards = [
        ("ChatGPT", "Answers anything, instantly.", TEAL),
        ("Notion AI", "Everything in one place.", PURPLE),
        ("Anki", "Flashcards built for memory.", ORANGE),
    ]
    gaps = [
        "No persistent learner model",
        "No spaced repetition",
        "No misconception diagnosis",
    ]
    xs = [1.0, 4.95, 8.9]
    card_w = 3.55
    top = 2.6
    for i, ((name, sub, accent), gap, x) in enumerate(zip(cards, gaps, xs)):
        rect(slide, x, top, card_w, 3.0, fill=SURFACE_2, line=accent)
        text(slide, name, x + 0.35, top + 0.32, card_w - 0.7, 0.3, size=18, color=TEXT, bold=True)
        text(slide, sub, x + 0.35, top + 0.82, card_w - 0.7, 0.3, size=11, color=MUTED)
        line(slide, x + 0.3, top + 1.43, x + card_w - 0.3, top + 1.43, accent, 1.2)
        rect(slide, x + 0.3, top + 1.78, card_w - 0.6, 0.75, fill=SURFACE_3, line=RED, radius=True)
        text(
            slide,
            gap,
            x + 0.35,
            top + 1.96,
            card_w - 0.7,
            0.45,
            size=10,
            color=RED,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    rect(slide, 0.85, 6.05, 12.0, 0.8, fill=BLUE, radius=True)
    text(
        slide,
        "Avenire is the only one that gets smarter about you — not just the material.",
        0.85,
        6.18,
        12.0,
        0.55,
        size=17,
        color=BG,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide)
    set_notes(
        slide,
        "Why Avenire wins the wedge:\n\nChatGPT answers but never models who you are as a learner. Notion AI stores your material but schedules no retrieval. Anki drills memory but cannot diagnose why an answer is wrong.\n\nAvenire is the only one that gets smarter about you, not just the material — persistent learner model, spaced repetition and misconception diagnosis all in one loop.",
    )
    return slide


def market_sizing(prs):
    slide = base_slide(prs, 1)
    title_block(
        slide,
        "Market sizing",
        "Rough math beats a blank slide.",
        "Order-of-magnitude — enough to plan a wedge, not a market study.",
    )
    rows = [
        ("TAM", "Global self-directed learners", "Learners studying on their own, worldwide.", "1B+", 1.0, BLUE),
        ("SAM", "India — college students using AI", "≈40M college students · ~15% already study with AI.", "≈6M", 0.35, TEAL),
        ("SOM", "Reachable in year one", "BITS · Symbiosis · waitlist → first engaged cohort.", "10K+", 0.12, ORANGE),
    ]
    for i, (tag, name, desc, value, frac, accent) in enumerate(rows):
        y = 2.55 + i * 1.05
        text(slide, tag, 0.95, y + 0.12, 1.0, 0.2, size=10, color=accent, bold=True, font=MONO)
        text(slide, name, 0.95, y + 0.38, 5.4, 0.25, size=15, color=TEXT, bold=True)
        text(slide, desc, 0.95, y + 0.72, 5.4, 0.3, size=10.5, color=MUTED)
        rect(slide, 6.9, y + 0.28, 4.6, 0.4, fill=SURFACE_3, radius=True)
        rect(slide, 6.9, y + 0.28, 4.6 * frac, 0.4, fill=accent, radius=True)
        text(slide, value, 11.7, y + 0.28, 1.0, 0.3, size=15, color=accent, bold=True, align=PP_ALIGN.RIGHT)
    text(
        slide,
        "Illustrative estimates — ballpark direction. Tighten with primary research before a formal round.",
        0.65,
        6.3,
        11.9,
        0.2,
        size=9.5,
        color=FAINT,
        font=MONO,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide)
    set_notes(
        slide,
        "Even rough numbers beat none.\n\nIndia has ~40M college students. Assuming ~15% already use AI tools to study gives an SAM of roughly 6M reachable students; global self-directed learners form the TAM; the SOM is the students reachable today from BITS, Symbiosis and the waitlist in year one.\n\nEvery figure on this slide is deliberately illustrative. Tighten each one with primary research before a formal round.",
    )
    return slide


def unit_economics_footnote(prs):
    slide = find_slide(prs, "BUSINESS MODEL")
    text(
        slide,
        "Unit economics (illustrative): ~₹200 COGS · ~₹600 paid ARPU · ~₹400 blended revenue per active student · month",
        0.85,
        6.68,
        11.7,
        0.18,
        size=8.5,
        color=FAINT,
        font=MONO,
    )
    set_notes(
        slide,
        "The current business model is straightforward: a free tier to reduce adoption friction, a ₹450 core plan, a ₹1,350 scholar plan, and an institutional tier coming next.\n\nUnit economics are hypothetical at this stage: roughly ₹200 cost of goods, ~₹600 paid ARPU and ~₹400 blended revenue per active student per month. These are hypotheses to validate with the first paying cohort.",
    )


def roadmap_dedupe(prs):
    slide = find_slide(prs, "ROADMAP")
    matches = []
    for shape in slide.shapes:
        collect_text_shapes(shape, "Personalized study plans", matches)
    if len(matches) > 1:
        for duplicate in matches[1:]:
            remove_shape(duplicate)


def team_cofounder(prs):
    slide = find_slide(prs, "WHY AVENIRE")
    text(slide, "Co-founder — Symbiosis University, Pune", 1.15, 4.85, 2.35, 0.24, size=12, color=MUTED)
    set_notes(
        slide,
        "The founder is inside the problem: a working product, built from the ground up, with a student perspective.\n\nThere is a co-founder studying at Symbiosis, Pune. Confirm her name and consent before presenting — add it here if she agrees, or dial the line back to keep this accurate.",
    )


def ask_heading(prs):
    slide = find_slide(prs, "THE ASK")
    replace_text(slide, "Our next milestone", "What this unlocks")
    set_notes(
        slide,
        "Why we are approaching the incubator: the product foundation is built, and the next step is validating the business.\n\nThe bar reads 'What this unlocks' — it is the direct output of the four asks, not a parallel roadmap: 50 waitlist → 100 active learners → first paying users → institutional pilots. We want to use the BITS ecosystem itself as one of the earliest testing and distribution environments.",
    )


def live_demo_note(prs):
    slide = find_slide(prs, "LIVE DEMO")
    rect(slide, 0.9, 2.3, 11.55, 3.5, fill=BG, line=TEAL)
    text(slide, "SCREEN SHARE", 1.2, 2.62, 2.0, 0.2, size=10, color=TEAL, bold=True, font=MONO)
    text(slide, "Live demo — screen share.", 1.1, 3.55, 11.1, 0.6, size=30, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    text(
        slide,
        "Double pendulum  ·  DNA double helix  ·  interactive concept models",
        1.1,
        4.45,
        11.1,
        0.25,
        size=13,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    line(slide, 1.1, 5.0, 12.3, 5.0, TEAL, 1.6)
    text(
        slide,
        "Open the real product here — let the room see the model respond.",
        1.1,
        5.22,
        11.1,
        0.24,
        size=10.5,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    set_notes(
        slide,
        "This is a screen-share moment, not a screenshot.\n\nOpen the real product — double pendulum, DNA double helix, interactive concept models — and let the student manipulate the model, not just look at a static image.\n\nThe slide is intentionally a note, so the deck never reads as a broken export.",
    )


def main():
    shutil.copy2(DECK, DECK.with_name("avenire_pitch_deck_backup_before_cleanup.pptx"))
    prs = Presentation(DECK)

    # 1) Reorder: move "AI has commoditized answers" (why now) right after Problem.
    why_now = find_slide(prs, "WHY NOW")
    move_slide(prs, why_now, 2)

    # 2) New competitive landscape immediately after it.
    comp = competitive_landscape(prs)
    move_slide(prs, comp, 3)

    # 3) New market sizing right after the Market slide.
    market = find_slide(prs, "MARKET")
    msize = market_sizing(prs)
    market_index = next(i for i, s in enumerate(prs.slides) if s is market)
    move_slide(prs, msize, market_index + 1)

    # 4) Content edits (targeted by slide identity).
    cover_hook(prs)
    rebuild_missing_piece_right(prs)
    technology_proof(prs)
    unit_economics_footnote(prs)
    roadmap_dedupe(prs)
    team_cofounder(prs)
    ask_heading(prs)
    live_demo_note(prs)

    renumber(prs)

    temp = DECK.with_name(".avenire_pitch_deck.cleanup.pptx")
    prs.save(temp)
    temp.replace(DECK)
    print(DECK)


if __name__ == "__main__":
    main()